from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    # パワポの新規作成
    prs = Presentation()

    # スライド追加用の関数
    def add_slide(title, content_points, notes=None):
        # レイアウト1: タイトルとコンテンツ
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル設定
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # 本文設定
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(24)
        
        # ノート設定
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        
        return slide

    # タイトルスライド追加用の関数
    def add_title_slide(title, subtitle, affiliation, name):
        # レイアウト0: タイトルスライド
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"{subtitle}\n\n{affiliation}\n{name}"

    # --- 1枚目: タイトル ---
    add_title_slide(
        "自閉症児の表情理解を深める\n拡張現実（AR）システムの開発と評価",
        "表情認識トレーニングにおけるARの有効性検証",
        "〇〇大学 〇〇学部 〇〇学科",
        "学籍番号: 〇〇〇〇　氏名: 〇〇 〇〇"
    )

    # --- 2枚目: 目次 ---
    add_slide(
        "目次",
        [
            "1. 研究背景",
            "2. 研究目的",
            "3. 提案システム",
            "4. 実験方法",
            "5. 実験結果",
            "6. 考察",
            "7. 結論"
        ]
    )

    # --- 3枚目: 研究背景 ---
    slide3 = add_slide(
        "1. 研究背景",
        [
            "自閉スペクトラム症（ASD）の課題",
            "・他者の感情や表情を認識・理解することが困難",
            "・社会的・情緒的コミュニケーションの阻害要因",
            "",
            "技術的アプローチ（AR）の可能性",
            "・ARは現実世界への不安を軽減し、理解を助ける",
            "・楽しみながら学べるツールとして期待"
        ],
        notes="参考文献: Griffiths et al. (2019)"
    )
    
    # 3枚目に注釈テキストボックスを追加
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(9)
    height = Inches(0.5)
    txBox = slide3.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "※ 自閉症児は表情認識（特に恐怖や怒り）の能力が低下していることが報告されている (Griffiths et al., 2019)"
    p.font.size = Pt(14)

    # --- 4枚目: 研究目的 ---
    add_slide(
        "2. 研究目的",
        [
            "本研究の目的",
            "1. ASD児と定型発達児（TD）の表情認識能力の比較",
            "2. ARシステムを用いたトレーニング効果の検証",
            "",
            "アプローチ",
            "・6つの基本感情（幸福、悲しみ、驚き、恐れ、嫌悪、怒り）",
            "・子どもが安全かつ親しみやすいARシステムの開発"
        ]
    )

    # --- 5枚目: 提案システム ---
    slide5 = add_slide(
        "3. 提案システム（FETA）",
        [
            "システム概要",
            "・UnityとMicrosoft Kinectを使用",
            "・コントローラー不要、ジェスチャー操作",
            "",
            "学習コンテンツ（3つの画像セット）",
            "1. グラフィック（イラスト）",
            "2. 同一人物の実写",
            "3. 異なる人物の実写"
        ]
    )
    # 図を貼る場所のダミー図形
    shape = slide5.shapes.add_shape(
        1, Inches(5.5), Inches(2), Inches(4), Inches(4)
    )
    shape.text = "【ここに図1, 図2, 図3を貼り付け】\n（グラフィック～実写の画像例）"
    shape.fill.solid()
    shape.fill.fore_color.rgb = 0xEEEEEE

    # --- 6枚目: 実験方法 ---
    slide6 = add_slide(
        "4. 実験方法",
        [
            "参加者",
            "・6歳〜9歳の児童 計30名（ASD児15名、TD児15名）",
            "",
            "実験手順（2日間）",
            "・1日目：評価 → トレーニング → 再評価",
            "・2日目：定着度の確認（再評価）",
            "",
            "評価指標",
            "・表情選択の正答数、保護者によるユーザビリティ評価"
        ]
    )
    # 図を貼る場所のダミー図形
    shape = slide6.shapes.add_shape(
        1, Inches(6), Inches(2), Inches(3.5), Inches(4.5)
    )
    shape.text = "【ここに表1（参加者特徴）と\n図4（実験ステップ）を貼り付け】"
    shape.fill.solid()
    shape.fill.fore_color.rgb = 0xEEEEEE

    # --- 7枚目: 実験結果1 ---
    slide7 = add_slide(
        "5. 実験結果①（全体傾向）",
        [
            "グループ間の比較",
            "・全体として定型発達児（TD）の方が正答率が高い",
            "",
            "トレーニング効果",
            "・両グループとも、1日目より2日目の成績が向上",
            "・AR介入による学習効果が確認された"
        ]
    )
    # 図を貼る場所のダミー図形
    shape = slide7.shapes.add_shape(
        1, Inches(1), Inches(4), Inches(8), Inches(3)
    )
    shape.text = "【ここに図12を貼り付け】\n（1日目と2日目の成績比較グラフ）\n※この発表で最も重要なデータです"
    shape.fill.solid()
    shape.fill.fore_color.rgb = 0xEEEEEE

    # --- 8枚目: 実験結果2 ---
    slide8 = add_slide(
        "5. 実験結果②（感情別・評価）",
        [
            "感情別の特徴",
            "・得意：「幸福」「悲しみ」は両グループとも高正答率",
            "・苦手：「嫌悪」「怒り」は正答率が低い（特にASD児）",
            "",
            "ユーザビリティ（保護者評価）",
            "・「安全性」「楽しさ」「使いやすさ」で高評価",
            "・子どもたちが意欲的に取り組めた"
        ]
    )
    # 図を貼る場所のダミー図形
    shape = slide8.shapes.add_shape(
        1, Inches(6.5), Inches(2), Inches(3), Inches(4.5)
    )
    shape.text = "【ここに図11（感情別反応）と\n表9（保護者評価）を貼り付け】"
    shape.fill.solid()
    shape.fill.fore_color.rgb = 0xEEEEEE

    # --- 9枚目: 考察 ---
    add_slide(
        "6. 考察",
        [
            "ARの有効性",
            "・視覚・聴覚刺激がASD児の学習意欲を向上させた",
            "・苦手な感情（嫌悪・怒り）も介入により改善傾向",
            "",
            "先行研究との関連",
            "・ASD児の表情認識の困難さは既存研究と一致",
            "・適切な支援ツール（AR）によるスキル向上の可能性を実証"
        ],
        notes="参考文献: Drimalla et al. (2021), Chen et al. (2015)"
    )

    # --- 10枚目: 結論 ---
    add_slide(
        "7. 結論",
        [
            "結論",
            "・開発したARシステムは、ASD児の表情理解および社会的スキルの向上に寄与する",
            "・安全で親しみやすい学習環境を実現",
            "",
            "今後の課題",
            "・サンプルサイズの拡大（現在は30名）",
            "・より長期的な効果の検証"
        ]
    )

    # --- 11枚目: 参考文献 ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "主要参考文献"
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    refs = [
        "[1] Griffiths, S. et al.: Impaired Recognition of Basic Emotions from Facial Expressions in Young People with Autism Spectrum Disorder. J. Autism Dev. Disord., 49, 2768–2778, 2019.",
        "[2] McMahon, D.D. et al.: Effects of Digital Navigation Aids on Adults with Intellectual Disabilities. J. Spec. Educ. Technol., 30, 157–165, 2015.",
        "[3] Drimalla, H. et al.: Facial expression mimicry and recognition in autism. Mol. Autism, 12, 27, 2021.",
        "[4] Azuma, R.T.: A Survey of Augmented Reality. Presence, 6, 355–385, 1997."
    ]
    
    for ref in refs:
        p = tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(16)
        p.space_after = Pt(10)

    # 保存
    prs.save('Thesis_Presentation_AR_ASD.pptx')
    print("パワーポイントファイルを作成しました: Thesis_Presentation_AR_ASD.pptx")

if __name__ == "__main__":
    create_presentation()